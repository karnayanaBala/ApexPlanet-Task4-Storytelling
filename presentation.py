from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()

# Slide 1 - Title
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
slide1.shapes.title.text = "Data Analytics Internship"
slide1.placeholders[1].text = "ApexPlanet Software Pvt. Ltd.\nKarnayana Bala\nSummer Internship 2026"

# Slide 2 - Objective
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Project Objective"
slide2.placeholders[1].text = """• Analyze real-world Sales Dataset (1000 records)
• Identify key business insights
• Perform Statistical Hypothesis Testing
• Build Interactive Dashboard
• Present findings to stakeholders"""

# Slide 3 - Key Findings
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Key Findings"
slide3.placeholders[1].text = """• Total Revenue: ₹139.40 Million
• Electronics is #1 Category (₹50.77M)
• Top Customer: Customer_335 (₹16.8L)
• Bengaluru has highest Avg Order Value
• Unit Price & Sales: Strong Correlation (0.69)
• Male vs Female sales: No significant difference"""

# Slide 4 - Hypothesis Testing
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Hypothesis Testing Results"
slide4.placeholders[1].text = """Test 1: Male vs Female Sales
• T-statistic: 0.6820
• P-value: 0.4954
• Result: No significant difference

Test 2: Electronics vs Other Categories
• T-statistic: 0.8293
• P-value: 0.4071
• Result: No significant difference"""

# Slide 5 - Recommendations
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = "Business Recommendations"
slide5.placeholders[1].text = """• Focus more investment on Electronics category
• Target Bengaluru & Pune markets (highest AOV)
• Higher priced products drive more revenue
• Equal marketing strategy for Male & Female
• Expand product range in top performing cities"""

# Slide 6 - Conclusion
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
slide6.shapes.title.text = "Conclusion"
slide6.placeholders[1].text = """• Successfully analyzed 1000 sales records
• Identified key revenue drivers
• Statistical tests validated business insights
• Interactive Power BI Dashboard created
• Data-driven recommendations provided"""

prs.save('ApexPlanet_Presentation.pptx')
print("Presentation saved!")